#!/usr/bin/env python3
"""Convert Marp-compatible Markdown slide decks to editable PowerPoint files.

Usage:
    python3 scripts/md_to_pptx.py modules/01-cloud-fundamentals/slides.md
    python3 scripts/md_to_pptx.py --all   # Convert all modules

Generates editable .pptx files with native text, tables, and code blocks.
"""

import os
import re
import sys
import glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def parse_front_matter(content):
    """Extract Marp front matter values."""
    fm = {}
    match = re.match(r'^---\s*\n(.*?)\n---', content, re.DOTALL)
    if match:
        for line in match.group(1).split('\n'):
            if ':' in line:
                key, val = line.split(':', 1)
                fm[key.strip()] = val.strip().strip("'\"")
    return fm


def split_slides(content):
    """Split Markdown content into individual slides."""
    # Remove front matter
    content = re.sub(r'^---\s*\n.*?\n---\s*\n', '', content, count=1, flags=re.DOTALL)
    # Split on slide separator
    raw_slides = re.split(r'\n---\s*\n', content)
    return [s.strip() for s in raw_slides if s.strip()]


def extract_speaker_notes(slide_text):
    """Extract speaker notes from HTML comments."""
    notes = []
    for match in re.finditer(r'<!--\s*(.*?)\s*-->', slide_text, re.DOTALL):
        notes.append(match.group(1).strip())
    # Remove comments from slide text
    clean = re.sub(r'<!--.*?-->', '', slide_text, flags=re.DOTALL).strip()
    return clean, '\n'.join(notes)


def parse_table(text):
    """Parse a Markdown table into headers and rows."""
    lines = [l.strip() for l in text.strip().split('\n') if l.strip()]
    if len(lines) < 2:
        return None, None
    headers = [c.strip() for c in lines[0].split('|') if c.strip()]
    rows = []
    for line in lines[2:]:  # Skip separator line
        cells = [c.strip() for c in line.split('|') if c.strip()]
        if cells:
            rows.append(cells)
    return headers, rows


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    return slide


def add_section_slide(prs, title):
    """Add a section divider slide."""
    layout = prs.slide_layouts[2]  # Section Header
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    return slide


def add_content_slide(prs, title, body_elements):
    """Add a content slide with title and body elements."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    # Build body text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    first = True
    for elem in body_elements:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if elem.get('type') == 'bullet':
            p.text = elem['text']
            p.level = elem.get('level', 0)
            p.font.size = Pt(18)
        elif elem.get('type') == 'text':
            p.text = elem['text']
            p.font.size = Pt(16)
        elif elem.get('type') == 'bold':
            run = p.add_run()
            run.text = elem['text']
            run.font.bold = True
            run.font.size = Pt(18)
        elif elem.get('type') == 'quote':
            p.text = elem['text']
            p.font.size = Pt(14)
            p.font.italic = True
            p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)

    # Add title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    # Calculate table dimensions
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    table_width = Inches(9)
    table_height = Inches(0.4) * num_rows
    left = Inches(0.5)
    top = Inches(1.2)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height)
    table = table_shape.table

    # Style header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x23, 0x2F, 0x3E)

    # Fill data rows
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            if c_idx < num_cols:
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = cell_text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                # Alternate row colors
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF3, 0xF4)

    return slide


def add_code_slide(prs, title, code, language=""):
    """Add a slide with a code block."""
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)

    # Add title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    # Add code block with dark background
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(9), Inches(5)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    code_box.line.fill.background()

    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code.strip()
    p.font.name = 'Courier New'
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)

    return slide


def add_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


def process_slide_content(prs, slide_text, notes_text):
    """Process a single slide's Markdown content and add it to the presentation."""
    lines = slide_text.split('\n')
    title = ""
    body_elements = []
    tables = []
    code_blocks = []
    current_table = []
    in_table = False
    in_code = False
    code_lang = ""
    code_content = []

    for line in lines:
        # Detect headings
        if line.startswith('# ') and not line.startswith('## '):
            title = line[2:].strip()
            continue
        if line.startswith('## '):
            title = line[3:].strip()
            continue

        # Detect code blocks
        if line.startswith('```') and not in_code:
            in_code = True
            code_lang = line[3:].strip()
            code_content = []
            continue
        if line.startswith('```') and in_code:
            in_code = False
            code_blocks.append({'lang': code_lang, 'code': '\n'.join(code_content)})
            continue
        if in_code:
            code_content.append(line)
            continue

        # Detect tables
        if '|' in line and not line.startswith('>'):
            if re.match(r'^\|?\s*[-:]+', line):
                continue  # Skip separator
            current_table.append(line)
            in_table = True
            continue
        elif in_table:
            tables.append('\n'.join(current_table))
            current_table = []
            in_table = False

        # Detect bullets
        if re.match(r'^[-*]\s', line):
            text = re.sub(r'^[-*]\s+', '', line)
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  # Remove bold markers
            text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)  # Remove links
            body_elements.append({'type': 'bullet', 'text': text, 'level': 0})
            continue

        # Detect numbered items
        if re.match(r'^\d+\.\s', line):
            text = re.sub(r'^\d+\.\s+', '', line)
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
            text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
            body_elements.append({'type': 'bullet', 'text': text, 'level': 0})
            continue

        # Detect blockquotes
        if line.startswith('>'):
            text = line.lstrip('> ').strip()
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
            text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
            body_elements.append({'type': 'quote', 'text': text})
            continue

        # Regular text
        if line.strip():
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', line.strip())
            text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
            body_elements.append({'type': 'text', 'text': text})

    if in_table and current_table:
        tables.append('\n'.join(current_table))

    # Determine slide type and create it
    slide = None

    # Section divider (title only, no body)
    if title and not body_elements and not tables and not code_blocks:
        slide = add_section_slide(prs, title)
    # Table slide
    elif tables:
        headers, rows = parse_table(tables[0])
        if headers and rows:
            slide = add_table_slide(prs, title or "Comparison", headers, rows)
            # If there are also bullets, add them as a separate slide
            if body_elements:
                extra = add_content_slide(prs, title or "Details", body_elements)
                add_notes(extra, "")
    # Code slide
    elif code_blocks:
        slide = add_code_slide(prs, title or "Code Example", code_blocks[0]['code'], code_blocks[0].get('lang', ''))
    # Content slide with bullets
    elif body_elements:
        slide = add_content_slide(prs, title or "Content", body_elements)
    # Fallback: title only
    elif title:
        slide = add_section_slide(prs, title)

    if slide:
        add_notes(slide, notes_text)

    return slide


def convert_md_to_pptx(md_path):
    """Convert a Marp Markdown file to an editable PowerPoint file."""
    with open(md_path, 'r') as f:
        content = f.read()

    fm = parse_front_matter(content)
    slides_md = split_slides(content)

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    for slide_text in slides_md:
        clean_text, notes = extract_speaker_notes(slide_text)
        if clean_text:
            process_slide_content(prs, clean_text, notes)

    # Save
    output_path = md_path.replace('.md', '.pptx')
    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 scripts/md_to_pptx.py <slides.md> | --all")
        sys.exit(1)

    if sys.argv[1] == '--all':
        md_files = sorted(glob.glob('modules/*/slides.md'))
        print(f"Converting {len(md_files)} slide decks to PowerPoint...")
        for md_file in md_files:
            try:
                output = convert_md_to_pptx(md_file)
                print(f"  Created: {output}")
            except Exception as e:
                print(f"  FAILED: {md_file} - {e}")
        print("Done.")
    else:
        md_file = sys.argv[1]
        if not os.path.exists(md_file):
            print(f"File not found: {md_file}")
            sys.exit(1)
        output = convert_md_to_pptx(md_file)
        print(f"Created: {output}")


if __name__ == '__main__':
    main()
